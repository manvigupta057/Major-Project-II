from email.mime import image
import os
from typing import List, Dict, Any
import pandas as pd
from config import Config
from PIL import Image
import easyocr

class Document:
    def __init__(self, content: str, metadata: Dict[str, Any] = None):
        self.content = content
        self.metadata = metadata or {}
    
    def __repr__(self):
        return f"Document(content='{self.content[:50]}...')"

class DocumentProcessor:
    def __init__(self, chunk_size: int = None, chunk_overlap: int = None):
        self.chunk_size = chunk_size or Config.CHUNK_SIZE
        self.chunk_overlap = chunk_overlap or Config.CHUNK_OVERLAP
    
    def create_chunks(self, text, metadata=None):
        chunks = []
        start = 0
        text_length = len(text)
        
        while start < text_length:
            end = start + self.chunk_size
            chunk_text = text[start:end]
            
            chunk_metadata = metadata.copy() if metadata else {}
            chunk_metadata['chunk_start'] = start
            chunk_metadata['chunk_end'] = end
            
            chunks.append(Document(content=chunk_text, metadata=chunk_metadata))
            
            start = end - self.chunk_overlap
            
            if start <= end - self.chunk_size:
                start = end
        
        return chunks

    def load_file(self, file_path, **kwargs):
        ext = os.path.splitext(file_path)[1].lower()
        
        if ext == '.csv':
            return self._load_csv(file_path, **kwargs)
        elif ext == '.pdf':
            return self._load_pdf(file_path)
        elif ext in ['.docx', '.doc']:
            return self._load_docx(file_path)
        elif ext in ['.pptx', '.ppt']:
            return self._load_pptx(file_path)
        elif ext in ['.jpg', '.jpeg', '.png', '.bmp', '.gif']:
            return self._load_image(file_path)
        elif ext in ['.xlsx', '.xls']:
            return self._load_excel(file_path)
        elif ext in ['.txt', '.md']:
            return self._load_text(file_path)
        else:
            print(f"Unsupported file type: {ext}")
            raise ValueError(f"Unsupported file type: {ext}")

    def _load_csv(self, file_path, text_columns=None, metadata_columns=None):
        df = pd.read_csv(file_path)
        return self._process_dataframe(df, file_path, text_columns, metadata_columns)

    def _load_excel(self, file_path, text_columns=None, metadata_columns=None):
        df = pd.read_excel(file_path)
        return self._process_dataframe(df, file_path, text_columns, metadata_columns)

    def _process_dataframe(self, df, file_path, text_columns=None, metadata_columns=None):
        loaded_data = []
        for idx, row in df.iterrows():
            if text_columns:
                text_parts = [f"{col}: {row[col]}" for col in text_columns if col in df.columns]
            else:
                text_parts = [f"{col}: {val}" for col, val in row.items()]
            
            text = "\n".join(text_parts)
            
            metadata = {'source': file_path, 'row_index': idx}
            if metadata_columns:
                for col in metadata_columns:
                    if col in df.columns:
                        metadata[col] = row[col]
            else:
                # Convert all to string to be safe for vector store compatibility
                safe_metadata = {k: str(v) for k, v in row.to_dict().items()}
                metadata.update(safe_metadata)
            
            loaded_data.append({'content': text, 'metadata': metadata})
        return loaded_data

    def _load_pdf(self, file_path):
        import pypdf
        loaded_data = []
        reader = pypdf.PdfReader(file_path)
        for i, page in enumerate(reader.pages):
            text = page.extract_text()
            if text:
                metadata = {'source': file_path, 'page_number': i + 1}
                loaded_data.append({'content': text, 'metadata': metadata})
        return loaded_data

    def _load_pptx(self, file_path):
        from pptx import Presentation
        loaded_data = []
        prs = Presentation(file_path)
        for i, slide in enumerate(prs.slides):
            text_parts = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_parts.append(shape.text)
            
            text = "\n".join(text_parts).strip()
            if text:
                metadata = {'source': file_path, 'slide_number': i + 1}
                loaded_data.append({'content': text, 'metadata': metadata})
        return loaded_data

    def _load_docx(self, file_path):
        import docx
        doc = docx.Document(file_path)
        text = "\n".join([para.text for para in doc.paragraphs if para.text.strip()])
        return [{'content': text, 'metadata': {'source': file_path}}]

    def _load_text(self, file_path):
        with open(file_path, 'r', encoding='utf-8') as f:
            text = f.read()
        return [{'content': text, 'metadata': {'source': file_path}}]
    
    def _load_image(self, file_path):
        import numpy as np
        
        reader = easyocr.Reader(['en'])
        
        img = Image.open(file_path)
        img_array = np.array(img)
        
        results = reader.readtext(img_array)
        
        text = "\n".join([item[1] for item in results])
        
        metadata = {'source': file_path, 'type': 'image'}
        
        return [{'content': text, 'metadata': metadata}]